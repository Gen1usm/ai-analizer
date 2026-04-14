from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

import pytesseract
from docx import Document
from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pypdf import PdfReader

SUPPORTED_EXTENSIONS = {
    ".docx",
    ".pdf",
    ".pptx",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".txt",
}


@dataclass
class FileParserError(Exception):
    message: str
    file_name: str = ""

    def __str__(self) -> str:
        return self.message


def extract_text_from_file(uploaded_file) -> str:
    extension = Path(uploaded_file.name).suffix.lower()
    file_name = uploaded_file.name

    if extension not in SUPPORTED_EXTENSIONS:
        allowed = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise FileParserError(
            f"Формат файла {extension or 'без расширения'} не поддерживается. "
            f"Разрешены: {allowed}.",
            file_name=file_name,
        )

    uploaded_file.seek(0)
    file_bytes = uploaded_file.read()
    uploaded_file.seek(0)

    if not file_bytes:
        raise FileParserError("Загруженный файл пустой.", file_name=file_name)

    parser_map = {
        ".txt": _parse_txt,
        ".docx": _parse_docx,
        ".pdf": _parse_pdf,
        ".pptx": _parse_pptx,
        ".png": _parse_image,
        ".jpg": _parse_image,
        ".jpeg": _parse_image,
        ".webp": _parse_image,
    }

    try:
        text = parser_map[extension](file_bytes)
    except FileParserError as exc:
        if not exc.file_name:
            exc.file_name = file_name
        raise
    except Exception as exc:
        raise FileParserError(
            f"Не удалось обработать файл {file_name}: {exc}",
            file_name=file_name,
        ) from exc

    normalized_text = text.strip()
    if not normalized_text:
        raise FileParserError(
            f"Из файла {file_name} не удалось извлечь текст.",
            file_name=file_name,
        )

    return normalized_text


def _parse_txt(file_bytes: bytes) -> str:
    encodings = ("utf-8", "utf-8-sig", "cp1251", "latin-1")
    for encoding in encodings:
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise FileParserError(
        "Текстовый файл не удалось декодировать в поддерживаемой кодировке."
    )


def _parse_docx(file_bytes: bytes) -> str:
    document = Document(BytesIO(file_bytes))
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs]
    return "\n".join(part for part in paragraphs if part)


def _parse_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(BytesIO(file_bytes))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(page.strip() for page in pages if page.strip())


def _parse_pptx(file_bytes: bytes) -> str:
    presentation = Presentation(BytesIO(file_bytes))
    chunks = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text.strip())
    return "\n".join(chunk for chunk in chunks if chunk)


def _parse_image(file_bytes: bytes) -> str:
    try:
        image = Image.open(BytesIO(file_bytes))
    except UnidentifiedImageError as exc:
        raise FileParserError(
            "Изображение повреждено или имеет неподдерживаемую структуру."
        ) from exc

    try:
        text = pytesseract.image_to_string(image, lang="eng+rus")
    except pytesseract.TesseractNotFoundError as exc:
        raise FileParserError(
            "Для распознавания текста на изображениях нужен установленный Tesseract OCR."
        ) from exc
    except pytesseract.TesseractError as exc:
        raise FileParserError(f"OCR не смог обработать изображение: {exc}") from exc

    return text
